"""
Document processing pipeline for RAG
Handles PDF, DOCX, PPTX, TXT extraction and cleaning
"""

import os
import re
from typing import List, Dict
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation
import logging

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Process various document formats and extract text"""
    
    def __init__(self):
        self.supported_formats = {'.pdf', '.docx', '.pptx', '.txt'}
    
    def extract_text(self, file_path: str, file_type: str = None) -> str:
        """Extract text from document based on file extension"""
        path = Path(file_path)
        
        if not file_type:
            file_type = path.suffix.lower()
        
        if file_type not in self.supported_formats:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        try:
            if file_type == '.pdf':
                return self._extract_pdf(file_path)
            elif file_type == '.docx':
                return self._extract_docx(file_path)
            elif file_type == '.pptx':
                return self._extract_pptx(file_path)
            elif file_type == '.txt':
                return self._extract_txt(file_path)
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            raise
    
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        text = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text.append(page.extract_text())
        return '\n'.join(text)
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        doc = Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs]
        return '\n'.join(paragraphs)
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    
    def _extract_txt(self, file_path: str) -> str:
        """Extract text from TXT"""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    def clean_text(self, text: str) -> str:
        """Clean extracted text - remove formatting artifacts, charts, etc."""
        # Remove chart data patterns (e.g., "0 10 20 30 40Percentage...")
        text = re.sub(r'\d+\s+\d+\s+\d+\s+\d+\s+\d+[A-Za-z]', '', text)
        # Remove lines that are mostly numbers (chart axes)
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            # Skip lines that are mostly numbers/spaces (chart data)
            if re.match(r'^[\d\s\.\%]+$', line.strip()) and len(line.strip()) > 5:
                continue
            # Skip very short lines that are just numbers
            if len(line.strip()) < 3 and line.strip().isdigit():
                continue
            cleaned_lines.append(line)
        text = '\n'.join(cleaned_lines)
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove special characters but keep punctuation and important symbols
        text = re.sub(r'[^\w\s\.\,\;\:\!\?\-\(\)\$\%\@\#\&]', '', text)
        # Remove multiple newlines
        text = re.sub(r'\n{3,}', '\n\n', text)
        # Remove standalone numbers that are likely chart artifacts
        text = re.sub(r'\s+\d+\s+', ' ', text)
        return text.strip()
    
    def chunk_text(
        self, 
        text: str, 
        chunk_size: int = 1000, 
        chunk_overlap: int = 200
    ) -> List[Dict[str, any]]:
        """Split text into overlapping chunks"""
        chunks = []
        words = text.split()
        
        if len(words) <= chunk_size:
            return [{
                'text': text,
                'start': 0,
                'end': len(text)
            }]
        
        start = 0
        while start < len(words):
            end = min(start + chunk_size, len(words))
            chunk_text = ' '.join(words[start:end])
            
            chunks.append({
                'text': chunk_text,
                'start': start,
                'end': end
            })
            
            start += (chunk_size - chunk_overlap)
        
        return chunks

